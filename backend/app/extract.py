import io

IMAGE_EXT = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "heic", "heif", "tiff"}


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_image(filename: str, content_type: str | None = None) -> bool:
    return (content_type or "").startswith("image/") or _ext(filename) in IMAGE_EXT


def extract_text(filename: str, data: bytes, content_type: str | None = None) -> str:
    ext = _ext(filename)
    ct = content_type or ""

    if ext == "pdf" or ct.endswith("pdf"):
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            return "\n\n".join((p.extract_text() or "") for p in reader.pages).strip()
        except Exception:
            return ""

    if ext == "docx" or "wordprocessingml" in ct:
        try:
            import docx

            doc = docx.Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts).strip()
        except Exception:
            return ""

    if ext == "pptx" or "presentationml" in ct:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts).strip()
        except Exception:
            return ""

    if is_image(filename, content_type):
        return ""

    return data.decode("utf-8", errors="ignore").strip()
